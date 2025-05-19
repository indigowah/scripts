import os
import sys
import openai
import docx
import pptx
import tempfile
import subprocess
import shutil
import threading
from concurrent.futures import ThreadPoolExecutor
from pptx import Presentation
from docx import Document
from pathlib import Path
from PyPDF2 import PdfReader

# Set your OpenAI API key here or use an environment variable
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
openai.api_key = OPENAI_API_KEY

def docx_to_markdown(docx_path):
    with tempfile.NamedTemporaryFile(suffix=".md", delete=False) as tmp_md:
        tmp_md_path = tmp_md.name
    try:
        subprocess.run(
            ["pandoc", str(docx_path), "-f", "docx", "-t", "markdown", "-o", tmp_md_path],
            check=True
        )
        with open(tmp_md_path, "r", encoding="utf-8") as f:
            md_text = f.read()
    finally:
        os.remove(tmp_md_path)
    return md_text

def pptx_to_markdown(pptx_path):
    prs = Presentation(pptx_path)
    md_lines = []
    for i, slide in enumerate(prs.slides):
        md_lines.append(f"# Slide {i+1}")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    md_lines.append(text)
        md_lines.append("")
    return "\n".join(md_lines)

def pdf_to_markdown(pdf_path):
    # Extract text from PDF
    reader = PdfReader(pdf_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    # Use OpenAI to convert to Markdown
    prompt = (
        "Convert the following PDF text to Obsidian-flavored Markdown. "
        "Preserve headings, lists, and formatting as much as possible.\n\n"
        f"{text}"
    )
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=4096,
        temperature=0.2,
    )
    return response.choices[0].message.content.strip()

def save_markdown(md_text, output_path):
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(md_text)

def convert_file(file_path, output_dir):
    ext = file_path.suffix.lower()
    if ext == ".docx":
        md = docx_to_markdown(file_path)
    elif ext == ".pptx":
        md = pptx_to_markdown(file_path)
    elif ext == ".pdf":
        md = pdf_to_markdown(file_path)
    else:
        return
    output_path = output_dir / (file_path.stem + ".md")
    save_markdown(md, output_path)
    print(f"Converted: {file_path} -> {output_path}")

def replicate_folder_structure(input_dir, output_dir):
    input_dir = Path(input_dir)
    output_dir = Path(output_dir)
    for root, dirs, _ in os.walk(input_dir):
        for d in dirs:
            src_path = Path(root) / d
            rel_path = src_path.relative_to(input_dir)
            dest_path = output_dir / rel_path
            dest_path.mkdir(parents=True, exist_ok=True)

def crawl_and_convert(input_dir, output_dir):
    replicate_folder_structure(input_dir, output_dir)
    input_dir = Path(input_dir)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    # Step 1: Count convertible files
    convertible_exts = [".docx", ".pptx", ".pdf"]
    all_files = []
    for root, _, files in os.walk(input_dir):
        for file in files:
            file_path = Path(root) / file
            if file_path.suffix.lower() in convertible_exts:
                all_files.append(file_path)
    total_files = len(all_files)
    if total_files == 0:
        print("No convertible files found.")
        return

    # Step 2: Convert files in parallel and show progress
    converted = 0
    converted_lock = threading.Lock()

    def convert_and_report(file_path):
        nonlocal converted
        rel_path = file_path.relative_to(input_dir)
        out_subdir = output_dir / rel_path.parent
        out_subdir.mkdir(parents=True, exist_ok=True)
        convert_file(file_path, out_subdir)
        with converted_lock:
            converted += 1
            print(f"{converted}/{total_files} files converted")

    with ThreadPoolExecutor() as executor:
        executor.map(convert_and_report, all_files)

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python marching_converter.py <input_dir> <output_dir>")
        sys.exit(1)
    crawl_and_convert(sys.argv[1], sys.argv[2])